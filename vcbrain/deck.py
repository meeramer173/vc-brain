"""Pitch-deck (.pptx) ingestion for the inbound application.

Design rule (same as everywhere): the LLM only works at the edge. Here that
edge is `intelligence.extract_deck`, which turns slide text into structured,
memo-shaped claims. This module does the deterministic parts:
  - `extract_text`  : pull readable text out of a .pptx (python-pptx, no
                      rendering, no vision) — shapes, tables, and speaker notes.
  - `verify_against_ledger` : a DETERMINISTIC cross-check of the deck's
                      self-reported metrics against what we INDEPENDENTLY
                      sourced (today: GitHub stars). This is the honesty
                      payoff — an inflated deck ("12k stars") is caught against
                      the real number in the ledger, no LLM involved.

A deck asserts nothing on its own: its numbers are *claims* to be verified, and
they never feed the deterministic Founder Score.
"""

import io
import sqlite3

from . import ledger

MAX_TEXT = 12000  # plenty for a deck; guards a pathological file


def extract_text(data: bytes) -> str:
    """Slide-by-slide text from a .pptx (bytes). Includes shape text, tables,
    and speaker notes. Raises on a non-pptx / corrupt file — the caller
    degrades gracefully."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    blocks: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                t = shape.text_frame.text.strip()
                if t:
                    parts.append(t)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
        if getattr(slide, "has_notes_slide", False):
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                parts.append("[speaker notes] " + note)
        if parts:
            blocks.append(f"--- Slide {i} ---\n" + "\n".join(parts))
    return "\n\n".join(blocks)[:MAX_TEXT]


def bare_handle(value: str) -> str:
    """Reduce a handle that may arrive as a URL or @handle to the bare
    username — the LLM often returns 'github.com/novareeves' where we need
    'novareeves' for entity resolution to merge onto the sourced profile."""
    s = str(value).strip().rstrip("/")
    if "/" in s:
        s = s.rsplit("/", 1)[-1]
    return s.lstrip("@")


def _ledger_github_stars(conn: sqlite3.Connection, entity_id: int) -> int | None:
    """Total stars we INDEPENDENTLY sourced for this founder from GitHub.
    None when we have no GitHub signal at all (can't verify)."""
    total, seen = 0, False
    for e in ledger.events_for(conn, entity_id):
        if e["source"] == "github":
            s = e["payload"].get("stars") or 0
            if s:
                total += s
                seen = True
    return total if seen else None


def verify_against_ledger(
    conn: sqlite3.Connection, entity_id: int, claimed: dict
) -> list[dict]:
    """Deterministic deck-claim vs independent-evidence check. Returns one row
    per claimed metric: {metric, claimed, known, status, note}. Status is
    'corroborated' | 'overstated' | 'contradicted' | 'unverifiable'."""
    rows: list[dict] = []

    cs = claimed.get("github_stars")
    if cs is not None:
        known = _ledger_github_stars(conn, entity_id)
        if not known:
            rows.append({"metric": "GitHub stars", "claimed": cs, "known": None,
                         "status": "unverifiable",
                         "note": "no independent GitHub signal in our ledger"})
        else:
            ratio = cs / max(known, 1)
            if ratio <= 1.5:
                status, note = "corroborated", f"ledger shows {known} across sourced repos"
            elif ratio >= 3:
                status, note = "contradicted", f"ledger shows only {known} across sourced repos"
            else:
                status, note = "overstated", f"ledger shows {known} across sourced repos"
            rows.append({"metric": "GitHub stars", "claimed": cs, "known": known,
                         "status": status, "note": note})

    for key, label in (("users", "Active users"), ("arr_usd", "ARR (USD)")):
        v = claimed.get(key)
        if v is not None:
            rows.append({"metric": label, "claimed": v, "known": None,
                         "status": "unverifiable",
                         "note": "self-reported; no independent source ingested"})
    return rows
