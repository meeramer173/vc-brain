"""Generate synthetic sample pitch decks (.pptx) for the inbound-application demo.

The brief permits synthetic decks. These two are built to exercise the deck
ingestion + the deterministic deck-vs-ledger Trust check:

  • DataForge AI (github novareeves) — INFLATED: claims 12,000 stars while the
    ledger holds ~41 for that founder → the check flags "contradicted".
  • NimbusRL (github sofia-nimbus)   — HONEST: claims ~3,300 stars, matching the
    2,400 + 900 the ledger already sourced → "corroborated".

Both founder github handles match seeded founders, so uploading the deck
resolves onto the existing profile ("we already knew you") and the claimed
traction is checked against what we independently sourced.

Usage:  python make_sample_decks.py   # writes data/sample_decks/*.pptx
"""

from pathlib import Path

from pptx import Presentation

OUT = Path(__file__).resolve().parent / "data" / "sample_decks"

DECKS = {
    "dataforge_inflated.pptx": [
        ("DataForge AI", ["AI data pipelines that just work",
                          "Nova Reeves · Founder · github.com/novareeves"]),
        ("Problem", ["Teams waste weeks building unreliable AI data pipelines",
                     "Existing tools break at production scale"]),
        ("Product", ["A managed, reliable AI data-pipeline toolkit",
                     "Drop-in, open-source core"]),
        ("Traction", ["12,000 GitHub stars", "50,000 active users",
                      "$1.2M ARR", "20% month-over-month growth"]),
        ("Market", ["AI data infrastructure — $30B TAM",
                    "Every AI team needs reliable pipelines"]),
        ("Team", ["Nova Reeves — ex-BigCo data-infra engineer"]),
        ("Ask", ["Raising a $100K first check to scale go-to-market"]),
    ],
    "nimbusrl_honest.pptx": [
        ("NimbusRL", ["Open-source RL training & serving infra for AI teams",
                      "Sofia Almeida · Founder · github.com/sofia-nimbus"]),
        ("Problem", ["AI teams can't serve reinforcement-learning models at low latency"]),
        ("Product", ["Open-source RL infrastructure plus low-latency model serving"]),
        ("Traction", ["3,300 GitHub stars across our repos",
                      "Published research paper", "Strong Show HN launch"]),
        ("Market", ["AI infrastructure — expanding fast; RL serving is a real need"]),
        ("Team", ["Sofia Almeida — RL researcher and open-source builder"]),
        ("Ask", ["Raising a $100K first check"]),
    ],
}


def _make(path: Path, slides: list) -> None:
    prs = Presentation()
    for i, (heading, bullets) in enumerate(slides):
        if i == 0:  # title slide
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = heading
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = "\n".join(bullets)
        else:       # title + content
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = heading
            tf = slide.placeholders[1].text_frame
            tf.text = bullets[0]
            for b in bullets[1:]:
                tf.add_paragraph().text = b
    prs.save(str(path))


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    for fname, slides in DECKS.items():
        _make(OUT / fname, slides)
        print(f"  wrote {OUT / fname}")


if __name__ == "__main__":
    main()
